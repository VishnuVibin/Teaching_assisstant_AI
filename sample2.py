import os
import re
import sys
import json
import subprocess
from dotenv import load_dotenv

load_dotenv()

try:
    import fitz 
except ImportError:
    print("PyMuPDF (fitz) is not installed. Installing PyMuPDF...")
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "pymupdf"])
        import fitz
    except Exception as e:
        print(f"Error installing PyMuPDF: {e}")
        sys.exit(1)

try:
    import pptx
except ImportError:
    print("python-pptx is not installed. Installing python-pptx...")
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        import pptx
        print("python-pptx installed successfully!")
    except Exception as e:
        print(f"Error installing python-pptx: {e}")
        sys.exit(1)

import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

API_KEY = os.getenv("OPENROUTER_API_KEY")
MODEL = "meta-llama/llama-3.1-70b-instruct"

def build_pptx(slides_data, output_path="teaching_slides.pptx"):
    prs = Presentation()

    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    for slide_info in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(11, 15, 25)
        
        if slide_info["type"] == "title":

            left_accent = slide.shapes.add_shape(
                1, 
                Inches(0), Inches(0), Inches(0.4), Inches(7.5)
            )
            left_accent.fill.solid()
            left_accent.fill.fore_color.rgb = RGBColor(124, 58, 237) 
            left_accent.line.fill.background()
            
            tb = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(2.0))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_info.get("title", "Lecture Presentation")
            p.font.name = "Arial"
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            sb = slide.shapes.add_textbox(Inches(1.2), Inches(4.2), Inches(11.0), Inches(1.5))
            sf = sb.text_frame
            sf.word_wrap = True
            p2 = sf.paragraphs[0]
            p2.text = slide_info.get("subtitle", "")
            p2.font.name = "Calibri"
            p2.font.size = Pt(20)
            p2.font.color.rgb = RGBColor(167, 139, 250) 
            
        else:
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.3), Inches(1.0))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_info.get("title", "Untitled Slide")
            p.font.name = "Arial"
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(6, 182, 212) 
            
            line = slide.shapes.add_shape(
                1, 
                Inches(1.0), Inches(1.4), Inches(11.3), Inches(0.04)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(124, 58, 237) 
            line.line.fill.background()

            cb = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.8))
            cf = cb.text_frame
            cf.word_wrap = True
            
            bullets = slide_info.get("bullets", [])
            for idx, bullet in enumerate(bullets):
                p = cf.add_paragraph() if idx > 0 else cf.paragraphs[0]
                p.text = bullet
                p.level = 0
                p.font.name = "Calibri"
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(241, 245, 249) 
                p.space_after = Pt(14)
                
    prs.save(output_path)
    return True

def extract_text(pdf_path):
    print("Extracting text from document...")
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text()
    doc.close()
    return text

def clean_and_parse_json(md_content):
    if not md_content or not md_content.strip():
        raise ValueError("Empty response received from LLM.")

    json_str = None
    json_block_match = re.search(r'```(?:json)?\s*([\s\S]*?)\s*```', md_content, re.IGNORECASE)
    if json_block_match:
        json_str = json_block_match.group(1).strip()
    else:
        array_match = re.search(r'(\[\s*\{[\s\S]*\}\s*\])', md_content)
        if array_match:
            json_str = array_match.group(1).strip()
        else:
            json_str = md_content.strip()

    try:
        return json.loads(json_str, strict=False)
    except Exception:
        pass

    def escape_raw_control_chars_in_strings(s):
        result = []
        in_string = False
        escaped = False
        for char in s:
            if char == '"' and not escaped:
                in_string = not in_string
                result.append(char)
            elif in_string:
                if char == '\n':
                    result.append('\\n')
                elif char == '\r':
                    result.append('\\r')
                elif char == '\t':
                    result.append('\\t')
                elif ord(char) < 32:
                    result.append(' ')
                else:
                    result.append(char)
            else:
                result.append(char)

            if char == '\\' and not escaped:
                escaped = True
            else:
                escaped = False
        return "".join(result)

    cleaned = escape_raw_control_chars_in_strings(json_str)

    try:
        return json.loads(cleaned, strict=False)
    except Exception:
        pass

    cleaned = re.sub(r',\s*([\]}])', r'\1', cleaned)
    cleaned = re.sub(r'\\(?!["\\/bfnrt]|u[0-9a-fA-F]{4})', r'\\\\', cleaned)

    try:
        return json.loads(cleaned, strict=False)
    except Exception:
        pass

    objects = re.findall(r'(\{\s*"type"\s*:\s*"[^"]+"[\s\S]*?\n\s*\})', json_str)
    parsed_slides = []
    for obj_str in objects:
        try:
            fixed_obj = escape_raw_control_chars_in_strings(obj_str)
            fixed_obj = re.sub(r',\s*([\]}])', r'\1', fixed_obj)
            fixed_obj = re.sub(r'\\(?!["\\/bfnrt]|u[0-9a-fA-F]{4})', r'\\\\', fixed_obj)
            slide_obj = json.loads(fixed_obj, strict=False)
            if isinstance(slide_obj, dict) and "type" in slide_obj:
                parsed_slides.append(slide_obj)
        except Exception:
            continue

    if parsed_slides:
        return parsed_slides

    return json.loads(json_str, strict=False)


def generate_slides(pdf_path, output_ppt_name="teaching_slides.pptx"):
    raw_text = extract_text(pdf_path)
    print(f"Extracted {len(raw_text)} characters.")
    
    context_text = raw_text[:40000]
    
    prompt = f"""You are a professional educational slide designer.
Analyze the following course material/chapter content and generate a highly structured 7-slide lecture layout.

Your output MUST be a valid JSON array matching the structure below, inside a ```json ``` code block. Do not write any other explanation or intro text outside the code block.
Ensure all string values are on a single line and contain NO unescaped raw newlines or control characters.

JSON Structure:
[
  {{
    "type": "title",
    "title": "Title of Lecture",
    "subtitle": "Course Name & Introduction"
  }},
  {{
    "type": "content",
    "title": "1. Lecture Learning Objectives",
    "bullets": [
      "Objective 1 description...",
      "Objective 2 description...",
      "Objective 3 description..."
    ]
  }},
  {{
    "type": "content",
    "title": "2. Core Concept: [Name]",
    "bullets": [
      "Definition and core idea...",
      "Analogy: [Real-world analogy to explain it]...",
      "Key teaching takeaway..."
    ]
  }},
  {{
    "type": "content",
    "title": "3. Technical Details: [Name]",
    "bullets": [
      "Technical details or implementation...",
      "How it works step-by-step...",
      "Common misconceptions to clarify..."
    ]
  }},
  {{
    "type": "content",
    "title": "4. Practical Application",
    "bullets": [
      "Real-world industry applications...",
      "Career relevance (e.g. roles that use this)...",
      "Discussion: [A prompt for student discussion]..."
    ]
  }},
  {{
    "type": "content",
    "title": "5. Active Classroom Exercise",
    "bullets": [
      "Exercise Title: [Think-pair-share / brief activity]...",
      "Step 1: Instructions for students...",
      "Step 2: Expected outcomes and talking points..."
    ]
  }},
  {{
    "type": "content",
    "title": "6. Summary & Check-on-Learning",
    "bullets": [
      "Summary of main points...",
      "Review Question 1...",
      "Review Question 2..."
    ]
  }}
]

Course Material Content:
{context_text}
"""

    print("Requesting slide layouts from OpenRouter API (meta-llama/llama-3.1-70b-instruct)...")

    try:
        response = requests.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {API_KEY}",
                "Content-Type": "application/json"
            },
            json={
                "model": MODEL,
                "temperature": 0.2,
                "messages": [
                    {
                        "role": "system",
                        "content": "You are a professional lecture slides creator. Return strict, valid JSON only without unescaped control characters inside string values."
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ]
            }
        )
        response.raise_for_status()
        md_content = response.json()["choices"][0]["message"]["content"]
    except Exception as e:
        print(f"API Error generating slides: {e}")
        return False
        
    print("Slide content generated. Parsing JSON structure...")
    
    try:
        slides_data = clean_and_parse_json(md_content)
    except Exception as e:
        print(f"Error parsing JSON slide content: {e}")
        print("Raw response preview:")
        print(md_content[:500])
        return False
        
    print("JSON parsed successfully. Compiling into PPTX slides...")
    
    try:
        build_pptx(slides_data, output_ppt_name)
        print(f"Success! PowerPoint presentation saved as '{output_ppt_name}' in the workspace.")
        return True
    except Exception as e:
        print(f"Error building PPTX slide deck: {e}")
        return False

def main():
    print("=" * 60)
    print("         SMART TA PORTAL - PPT LECTURE SLIDE DECK GENERATOR")
    print("=" * 60)
    
    pdf_path = input("Enter course PDF material path: ").strip()
    if not pdf_path or not os.path.exists(pdf_path):
        print(f"Error: File '{pdf_path}' not found.")
        return
        
    output_name = input("Enter output PPTX name [default: teaching_slides.pptx]: ").strip()
    if not output_name:
        output_name = "teaching_slides.pptx"
    if not output_name.lower().endswith(".pptx"):
        output_name += ".pptx"
        
    generate_slides(pdf_path, output_name)
if __name__ == "__main__":
    main()
